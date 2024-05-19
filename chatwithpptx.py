import streamlit as st
import os
from pptx import Presentation
from streamlit_extras.add_vertical_space import add_vertical_space
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import AzureOpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.llms import AzureOpenAI
from langchain.chains.question_answering import load_qa_chain
from langchain.callbacks import get_openai_callback
from dotenv import load_dotenv
import pickle

load_dotenv()

os.environ["AZURE_OPENAI_API_KEY"] = "your-azure-openai-api-key"
os.environ["AZURE_OPENAI_API_BASE"] = "your-azure-openai-api-base"
os.environ["AZURE_OPENAI_API_VERSION"] = "your-azure-openai-api-version"
os.environ["AZURE_OPENAI_DEPLOYMENT_NAME"] = "your-deployment-name"

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def my_function():
    st.write("Function called!")
    vector_list = []
    folder_path = '/home/ai/project/Langchain-Document-Chat/data'
    for root, dirs, filenames in os.walk(folder_path):
        text = ""
        for file in filenames:
            fileL = os.path.join(root, file)
            print(fileL)
            text += extract_text_from_pptx(fileL)

        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        chunks = text_splitter.split_text(text=text)
        embeddings = AzureOpenAIEmbeddings(deployment_name=os.environ["AZURE_OPENAI_DEPLOYMENT_NAME"])
        print(chunks)
        VectorStore = FAISS.from_texts(chunks, embedding=embeddings)
        vector_list.append(VectorStore)
        print(vector_list)
        store_name = "master.pickle"
        vector_store_path = os.path.join('vector_store', store_name)    
        with open(vector_store_path, 'wb') as file:
                pickle.dump(VectorStore, file)
                print('wow')   
    print("Successfully done")

def main():
    st.header("Private Chat Window - You can ask any query 💬")
    pptx = st.file_uploader("Upload your PPTX", type="pptx")

    if pptx is not None:
        text = extract_text_from_pptx(pptx)

        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        chunks = text_splitter.split_text(text=text)
    
        store_name = pptx.name[:-5] + ".pickle"
        vector_store_path = os.path.join('vector_store', store_name)
    else:
        store_name = "master.pickle"
        vector_store_path = os.path.join('vector_store', store_name) 

    if os.path.exists(vector_store_path):
        print("I already have it")
        with open(vector_store_path, 'rb') as file:
            VectorStore = pickle.load(file)
            print(VectorStore)
            print(type(VectorStore))
    else:
        embeddings = AzureOpenAIEmbeddings(deployment_name=os.environ["AZURE_OPENAI_DEPLOYMENT_NAME"])
        VectorStore = FAISS.from_texts(chunks, embedding=embeddings)
        with open(vector_store_path, 'wb') as file:
            pickle.dump(VectorStore, file)

    chat_history = st.session_state.get("chat_history", [])

    query = st.text_input("Ask questions about your PPTX file")

    if query:
        docs = VectorStore.similarity_search(query=query, k=3)
        llm = AzureOpenAI(deployment_name=os.environ["AZURE_OPENAI_DEPLOYMENT_NAME"], model_name='gpt-3.5-turbo')
        chain = load_qa_chain(llm=llm, chain_type='stuff')
        with get_openai_callback() as cb:
            response = chain.run(input_documents=docs, question=query)
            response_with_citations = f"{response}\n\nCitations:\n" + "\n".join([doc.metadata['source'] for doc in docs])
            chat_history.append({"user": query, "bot": response_with_citations})

        st.session_state.chat_history = chat_history

        for i, message in enumerate(chat_history):
            with st.chat_message("user" if i % 2 == 0 else "assistant"):
                st.markdown(f"**{message['user']}**")
                st.markdown(message['bot'])

if __name__ == "__main__":
    main()
